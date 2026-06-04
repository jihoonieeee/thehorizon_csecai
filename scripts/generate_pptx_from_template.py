#!/usr/bin/env python3
"""
generate_pptx_from_template.py

Generates a branded PPTX from the AI x Security template using slide content JSON.
Uses python-pptx to inherit the template's master slide, layouts, fonts, and colours.

Usage:
    python3 scripts/generate_pptx_from_template.py \
        --input  outputs/final/slide_deck_output.json \
        --output outputs/final/horizon_scan_deck.pptx \
        --template "templates/AI x Security (for AISP projection) (1).pptx"

If --template is omitted, falls back to the default template path.
"""

import sys
import json
import argparse
import copy
import io
import re
import zipfile
import warnings
from pathlib import Path

# Suppress python-pptx duplicate-name zip warnings (handled by cleanup step below)
warnings.filterwarnings("ignore", category=UserWarning, module="zipfile")

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from lxml import etree

# Global visualization index: visualization_id → spec. Populated in generate().
VIZ_INDEX = {}

# ── Defaults ──────────────────────────────────────────────────────────────────

SCRIPT_DIR   = Path(__file__).parent
PROJECT_ROOT = SCRIPT_DIR.parent
DEFAULT_TEMPLATE = PROJECT_ROOT / "templates" / "AI x Security (for AISP projection) (1).pptx"

# ── Colour palette (CSA/AISP brand) ──────────────────────────────────────────

NAVY    = RGBColor(0x00, 0x49, 0x87)
BLUE    = RGBColor(0x35, 0x83, 0xC9)
PURPLE  = RGBColor(0x9C, 0x62, 0xA7)
TEAL    = RGBColor(0x19, 0xBC, 0x9D)
AMBER   = RGBColor(0xFF, 0xAA, 0x22)
RED     = RGBColor(0xCC, 0x00, 0x33)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
GREY    = RGBColor(0x6B, 0x72, 0x80)
OFFWHITE = RGBColor(0xF4, 0xF6, 0xF9)

CATEGORY_COLORS = {
    "traditional_ai_threats": PURPLE,
    "llm_threats":            BLUE,
    "agentic_ai_threats":     TEAL,
    "ai_enabled_threats":     AMBER,
    "cross_category":         RED,
}

# ── Layout mapping ────────────────────────────────────────────────────────────
# (master_idx, layout_idx_within_master)
# Master 1: 0=Cover, 1=TitleContent, 2=TwoContent, 3=Comparison, 4=Divider, 5=Blank, 6=SectionHeader, 7=3_Title
# Master 2: 0=4_Title, 1=TitleContent, 2=TwoContent, 3=Comparison, 4=TitleOnly, 5=SectionHeader

LAYOUT_MAP = {
    "title":                    (0, 0),  # Master1 Cover Page
    # Overview
    "scope_methodology":        (0, 1),  # Master1 Title and Content
    "scope_timeframe":          (0, 1),
    "methodology":              (0, 1),
    "source_coverage":          (0, 1),
    "taxonomy_framework":       (0, 1),
    "landscape":                (1, 1),  # Master2 Title and Content
    "corpus_analytics":         (1, 1),
    "cross_cutting_trends":     (1, 1),
    "cross_cutting_signals":    (1, 1),
    "exec_overview":            (1, 1),
    # Per-category
    "section_divider":          (0, 4),  # Master1 Divider
    "category_viewpoint":       (1, 1),
    "category_technique_map":   (1, 1),
    "category_evidence":        (1, 1),
    "case_studies":             (1, 1),
    "category_analytics":       (1, 1),
    "category_outlook_gaps":    (1, 1),
    "category_analytics_outlook": (1, 1),
    "category_content":         (1, 1),
    # Synthesis
    "cross_category":           (1, 1),
    "maturity_assessment":      (1, 1),
    "governance_implications":  (1, 1),
    "recommendations":          (1, 1),
    "watchlist":                (0, 1),
    "watchlist_gaps":           (0, 1),
    "evidence_gaps_confidence": (0, 1),
    "outlook":                  (1, 1),
    "conclusion":               (1, 5),  # Master2 Section Header
    # Appendix
    "appendix":                 (0, 1),
    "appendix_evidence_index":  (0, 1),
    "appendix_analytics_tables":(0, 1),
    "appendix_taxonomy":        (0, 1),
}
DEFAULT_LAYOUT = (1, 1)

# ── Helpers ───────────────────────────────────────────────────────────────────

def get_layout(prs, master_idx, layout_idx):
    master = prs.slide_masters[min(master_idx, len(prs.slide_masters) - 1)]
    layout = master.slide_layouts[min(layout_idx, len(master.slide_layouts) - 1)]
    return layout


def set_ph_text(slide, idx, text, bold=False, size_pt=None, color=None, align=None):
    """Set the text of a placeholder by idx. Returns True if found."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        if bold:
                            run.font.bold = bold
                        if size_pt:
                            run.font.size = Pt(size_pt)
                        if color:
                            run.font.color.rgb = color
                    if align:
                        para.alignment = align
            return True
    return False


def clear_ph(slide, idx):
    """Clear a placeholder's text frame."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            if ph.has_text_frame:
                for p in ph.text_frame.paragraphs:
                    for r in p.runs:
                        r.text = ""
            return ph
    return None


def fill_body(slide, headline, bullets, evidence_callouts, citations, analytics_hints=None):
    """
    Fill the body placeholder (idx=1) with:
    - headline (bold, slightly larger)
    - bullet points (standard)
    - evidence callouts (indented, publisher bold + fact)
    - optional analytics hint line
    """
    ph = None
    for p in slide.placeholders:
        if p.placeholder_format.idx == 1:
            ph = p
            break
    if not ph or not ph.has_text_frame:
        return

    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True

    def add_para(tf, text, level=0, bold=False, italic=False, size_pt=None, color=None, align=None, space_before=0):
        """Add a paragraph with runs to the text frame."""
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn as _qn
        p = tf.add_paragraph()
        p.level = level
        if space_before:
            p.space_before = Pt(space_before)
        if align:
            p.alignment = align
        run = p.add_run()
        run.text = text
        if bold:
            run.font.bold = True
        if italic:
            run.font.italic = True
        if size_pt:
            run.font.size = Pt(size_pt)
        if color:
            run.font.color.rgb = color
        return p

    # Clear the first auto-paragraph
    first_para = tf.paragraphs[0]
    first_para.clear()

    # Headline
    if headline:
        p = first_para
        p.level = 0
        run = p.add_run()
        run.text = headline
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = NAVY

    # Bullets
    for bullet in (bullets or []):
        add_para(tf, bullet, level=0, size_pt=11)

    # Analytics hint line
    if analytics_hints:
        add_para(tf, "", level=0)
        add_para(tf, analytics_hints, level=0, italic=True, size_pt=9, color=GREY)

    # Evidence callouts
    if evidence_callouts:
        add_para(tf, "", level=0)
        add_para(tf, "KEY EVIDENCE", level=0, bold=True, size_pt=9, color=NAVY, space_before=4)
        for ev in (evidence_callouts or [])[:3]:
            publisher = (ev.get("publisher") or "").upper()
            key_fact  = (ev.get("key_fact") or "")[:150]
            # Publisher line
            p = tf.add_paragraph()
            p.level = 1
            r1 = p.add_run()
            r1.text = f"{publisher}  "
            r1.font.bold = True
            r1.font.size = Pt(8)
            r1.font.color.rgb = BLUE
            r2 = p.add_run()
            r2.text = key_fact
            r2.font.size = Pt(8)
            r2.font.color.rgb = DARK


def set_speaker_notes(slide, notes_text):
    """Write speaker notes to the slide's notes pane."""
    if not notes_text:
        return
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()
        tf.paragraphs[0].text = notes_text
    except Exception:
        pass


def add_accent_bar(slide, color, left=0, top=0, width=None, height=Inches(0.08)):
    """Add a thin horizontal colour bar (used under title in content slides)."""
    w = width or Inches(13.33)
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(int(left)), Emu(int(top)), Emu(int(w)), Emu(int(height))
    )
    shape.line.color.rgb = color
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    # Remove line border
    shape.line.width = Emu(0)
    return shape


def add_text_box(slide, text, left, top, width, height,
                 size_pt=12, bold=False, italic=False, color=None,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a freeform text box to the slide."""
    txBox = slide.shapes.add_textbox(
        Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


# ── Visualization rendering ───────────────────────────────────────────────────

CHART_SERIES_COLORS = [BLUE, PURPLE, TEAL, AMBER, NAVY, RED]


def normalize_chart_data(spec):
    """
    Normalize the heterogeneous chart_data formats into a canonical form:
      {"kind": "bar",   "categories": [...], "values": [...]}
      {"kind": "table", "columns": [...], "rows": [{"label", "values"}]}
      {"kind": "timeline", "events": [{"date", "label"}]}
    Returns None if the spec has no renderable data.
    """
    cd = spec.get("chart_data") or {}
    vtype = spec.get("visualization_type") or ""

    # items[] → bar
    if isinstance(cd.get("items"), list) and cd["items"]:
        cats = [(i.get("label") or i.get("key") or "") for i in cd["items"]]
        vals = [(i.get("count") or i.get("value") or 0) for i in cd["items"]]
        cats, vals = _drop_unknown(cats, vals)
        if cats:
            return {"kind": "bar", "categories": cats[:10], "values": vals[:10]}

    # stacked_bar {categories, stacks[{values}]} → summed bar
    if isinstance(cd.get("stacks"), list) and isinstance(cd.get("categories"), list):
        cats = cd["categories"]
        totals = [sum(s.get("values", [0])[i] if i < len(s.get("values", [])) else 0
                      for s in cd["stacks"]) for i in range(len(cats))]
        if any(totals):
            return {"kind": "bar", "categories": cats[:10], "values": totals[:10]}

    # stacked_bar {months, series[{values}]} → summed bar
    if isinstance(cd.get("series"), list) and isinstance(cd.get("months"), list):
        cats = cd["months"]
        totals = [sum(s.get("values", [0])[i] if i < len(s.get("values", [])) else 0
                      for s in cd["series"]) for i in range(len(cats))]
        if any(totals):
            return {"kind": "bar", "categories": cats[:12], "values": totals[:12]}

    # radar {axes, values} → bar
    if isinstance(cd.get("axes"), list) and isinstance(cd.get("values"), list):
        cats, vals = _drop_unknown(cd["axes"], cd["values"])
        if cats:
            return {"kind": "bar", "categories": cats[:10], "values": vals[:10]}

    # heatmap/matrix {columns, rows[{label, values}]} → table
    if isinstance(cd.get("rows"), list) and isinstance(cd.get("columns"), list) and cd["rows"]:
        return {"kind": "table", "columns": cd["columns"][:6],
                "rows": [{"label": r.get("label", ""), "values": (r.get("values") or [])[:6]}
                         for r in cd["rows"][:6]]}

    # timeline {events}
    if isinstance(cd.get("events"), list) and cd["events"]:
        return {"kind": "timeline",
                "events": [{"date": (e.get("date") or "")[:7], "label": e.get("title") or e.get("label") or ""}
                           for e in cd["events"][:8]]}

    return None


def _drop_unknown(cats, vals):
    """Drop 'unknown' category entries so charts show meaningful data only."""
    out_c, out_v = [], []
    for c, v in zip(cats, vals):
        if str(c).strip().lower() in ("unknown", "", "n/a", "none"):
            continue
        out_c.append(str(c).replace("_", " "))
        out_v.append(v)
    return out_c, out_v


def render_visual(slide, spec, x, y, w, h):
    """
    Render a visualization spec onto the slide as a native PPTX chart or table.
    Returns True if something was drawn.
    """
    norm = normalize_chart_data(spec)
    if not norm:
        return False

    title = spec.get("title", "")

    if norm["kind"] == "bar":
        chart_data = CategoryChartData()
        chart_data.categories = norm["categories"]
        chart_data.add_series(title or "Count", norm["values"])
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)),
            chart_data,
        )
        chart = gf.chart
        chart.has_legend = False
        chart.has_title = bool(title)
        if title:
            chart.chart_title.text_frame.text = title[:60]
            for run in chart.chart_title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = NAVY
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            plot.data_labels.font.size = Pt(8)
            for i, pt in enumerate(plot.series[0].points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = CHART_SERIES_COLORS[i % len(CHART_SERIES_COLORS)]
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)
        except Exception:
            pass
        return True

    if norm["kind"] == "table":
        cols = norm["columns"]; rows = norm["rows"]
        n_rows = len(rows) + 1
        n_cols = len(cols) + 1
        gf = slide.shapes.add_table(n_rows, n_cols, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
        tbl = gf.table
        tbl.cell(0, 0).text = ""
        for j, col in enumerate(cols):
            c = tbl.cell(0, j + 1)
            c.text = str(col).replace("_", " ")[:14]
            c.text_frame.paragraphs[0].runs[0].font.size = Pt(7)
            c.text_frame.paragraphs[0].runs[0].font.bold = True
        maxv = max([max(r["values"]) if r["values"] else 0 for r in rows] + [1])
        for i, r in enumerate(rows):
            lc = tbl.cell(i + 1, 0)
            lc.text = str(r["label"]).replace("_", " ")[:18]
            lc.text_frame.paragraphs[0].runs[0].font.size = Pt(7)
            for j in range(len(cols)):
                val = r["values"][j] if j < len(r["values"]) else 0
                cell = tbl.cell(i + 1, j + 1)
                cell.text = str(val) if val else ""
                if cell.text_frame.paragraphs[0].runs:
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(7)
                # Colour intensity by value
                t = (val / maxv) if maxv else 0
                if t > 0:
                    shade = int(0xF4 - (0xF4 - 0x35) * t)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(shade, shade, 0xC9)
        return True

    if norm["kind"] == "timeline":
        lines = [f"{e['date']}  {e['label'][:50]}" for e in norm["events"]]
        add_text_box(slide, "\n".join(lines),
                     Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)),
                     size_pt=9, color=DARK)
        return True

    return False


def pick_renderable_viz(slide_data):
    """Return the first visualization spec that has renderable data for this slide."""
    for vid in (slide_data.get("visualization_ids") or []):
        spec = VIZ_INDEX.get(vid)
        # Honour the JS thin-data guard: specs flagged insufficient_data
        # (<2 meaningful points or a single time bucket) are not charted.
        if spec and spec.get("insufficient_data"):
            continue
        if spec and normalize_chart_data(spec):
            return spec
    return None


# ── Slide builders ────────────────────────────────────────────────────────────

def build_title_slide(prs, layout, slide_data):
    slide = prs.slides.add_slide(layout)
    title    = slide_data.get("title", "AI Cyber Threat Horizon Scan")
    subtitle = slide_data.get("headline") or slide_data.get("core_message") or ""
    set_ph_text(slide, 0, title,    bold=False)
    set_ph_text(slide, 1, subtitle, bold=False, size_pt=16)
    set_speaker_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


def build_section_divider(prs, layout, slide_data):
    slide    = prs.slides.add_slide(layout)
    title    = slide_data.get("title", "")
    subtitle = slide_data.get("core_message") or slide_data.get("headline") or ""
    category = slide_data.get("category", "")
    accent   = CATEGORY_COLORS.get(category, BLUE)

    set_ph_text(slide, 0, title, bold=True, size_pt=36, color=WHITE)

    # Add core message below the title area as a text box
    if subtitle:
        add_text_box(
            slide, subtitle,
            left=Inches(0.8), top=Inches(4.8),
            width=Inches(11.7), height=Inches(1.2),
            size_pt=15, color=RGBColor(0xAA, 0xBB, 0xCC)
        )

    # Accent bar at top
    add_accent_bar(slide, accent, left=0, top=0, width=Inches(0.22), height=Inches(7.5))

    set_speaker_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


def build_conclusion_slide(prs, layout, slide_data):
    slide = prs.slides.add_slide(layout)
    title = slide_data.get("title", "Key Takeaways")
    set_ph_text(slide, 0, title, bold=True, size_pt=32)
    bullets = slide_data.get("bullets") or []
    # For Section Header (no body placeholder), add bullets as text box
    content = "\n".join(f"• {b}" for b in bullets[:5])
    if content:
        add_text_box(
            slide, content,
            left=Inches(0.8), top=Inches(2.5),
            width=Inches(11.7), height=Inches(4.0),
            size_pt=14, color=WHITE
        )
    set_speaker_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


def build_content_slide(prs, layout, slide_data):
    slide = prs.slides.add_slide(layout)

    title    = slide_data.get("title", "")
    headline = slide_data.get("headline", "")
    bullets  = slide_data.get("bullets") or []
    ev_calls = slide_data.get("evidence_callouts") or []
    citations = slide_data.get("citations") or []

    set_ph_text(slide, 0, title, bold=True)

    # If a renderable chart is assigned, lay out text on the left and chart on the right.
    spec = pick_renderable_viz(slide_data)
    if spec:
        # Narrow the body placeholder to the left ~55% so the chart has room.
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                try:
                    ph.width = Inches(7.0)
                except Exception:
                    pass
                break
        fill_body(slide, headline, bullets, ev_calls, citations[:2])
        drawn = render_visual(
            slide, spec,
            x=Inches(7.7), y=Inches(1.7), w=Inches(5.2), h=Inches(4.6)
        )
        if not drawn:
            pass
    else:
        fill_body(slide, headline, bullets, ev_calls, citations[:2])

    set_speaker_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


def build_source_coverage_slide(prs, layout, slide_data):
    slide = prs.slides.add_slide(layout)
    title = slide_data.get("title", "Source Coverage")
    set_ph_text(slide, 0, title, bold=True)

    agg     = slide_data.get("aggregates_summary") or {}
    total   = agg.get("total_sources", 0)
    cats    = agg.get("category_counts") or {}
    types   = agg.get("source_type_counts") or {}
    headline = f"{total} validated sources validated across {len(cats)} threat categories"
    bullets  = [
        f"Total sources: {total}",
        *[f"{k.replace('_',' ').title()}: {v} sources" for k,v in
          sorted(cats.items(), key=lambda x: -x[1])[:4]],
    ][:5]
    fill_body(slide, headline, bullets, [], [])
    set_speaker_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


def build_appendix_slide(prs, layout, slide_data):
    """Render an appendix slide: title + citation list in two columns."""
    slide = prs.slides.add_slide(layout)
    title = slide_data.get("title", "Appendix")
    set_ph_text(slide, 0, title, bold=True)

    citations = slide_data.get("citations") or []
    if citations:
        half = (len(citations) + 1) // 2
        left = citations[:half]
        right = citations[half:half * 2]
        col_w = Inches(6.1)
        add_text_box(slide, "\n".join(f"{i+1}. {c[:90]}" for i, c in enumerate(left)),
                     Inches(0.4), Inches(1.4), col_w, Inches(5.6), size_pt=7, color=DARK)
        if right:
            add_text_box(slide, "\n".join(f"{half+i+1}. {c[:90]}" for i, c in enumerate(right)),
                         Inches(6.7), Inches(1.4), col_w, Inches(5.6), size_pt=7, color=DARK)
    else:
        # No citations (e.g. taxonomy reference) — render headline + bullets
        fill_body(slide, slide_data.get("headline", ""), slide_data.get("bullets") or [], [], [])

    set_speaker_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


# ── Main dispatch ─────────────────────────────────────────────────────────────

BUILDER_MAP = {
    "title":                    build_title_slide,
    "section_divider":          build_section_divider,
    "conclusion":               build_conclusion_slide,
    "source_coverage":          build_source_coverage_slide,
    "appendix":                 build_appendix_slide,
    "appendix_evidence_index":  build_appendix_slide,
    "appendix_analytics_tables":build_appendix_slide,
    "appendix_taxonomy":        build_appendix_slide,
}

def build_slide(prs, slide_data):
    slide_type = slide_data.get("slide_type", "unknown")
    master_idx, layout_idx = LAYOUT_MAP.get(slide_type, DEFAULT_LAYOUT)
    layout  = get_layout(prs, master_idx, layout_idx)
    builder = BUILDER_MAP.get(slide_type, build_content_slide)
    return builder(prs, layout, slide_data)


# ── Zip cleanup — strip orphaned slide parts ──────────────────────────────────

def _collect_referenced_parts(zip_path: str) -> set:
    """
    Walk the relationship graph starting from presentation.xml to find all
    actually-referenced parts. Returns a set of normalised zip names to keep.
    """
    REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
    kept = set()

    def normalise(base: str, target: str) -> str:
        """Resolve a relative relationship target to an absolute zip path."""
        if target.startswith("/"):
            return target.lstrip("/")
        base_dir = base.rsplit("/", 1)[0] if "/" in base else ""
        parts = (base_dir + "/" + target).split("/")
        resolved = []
        for p in parts:
            if p == "..":
                if resolved:
                    resolved.pop()
            elif p and p != ".":
                resolved.append(p)
        return "/".join(resolved)

    def crawl(zf, partname: str):
        if partname in kept:
            return
        kept.add(partname)
        rels_path = ""
        base_dir, base_file = partname.rsplit("/", 1) if "/" in partname else ("", partname)
        rels_path = f"{base_dir}/_rels/{base_file}.rels" if base_dir else f"_rels/{base_file}.rels"
        if rels_path not in zf.namelist():
            return
        rels_xml = zf.read(rels_path).decode("utf-8")
        kept.add(rels_path)
        for m in re.finditer(r'Target="([^"]+)"', rels_xml):
            target = m.group(1)
            if target.startswith("http") or target.startswith("#"):
                continue
            resolved = normalise(partname, target)
            if resolved and resolved in zf.namelist():
                crawl(zf, resolved)

    with zipfile.ZipFile(zip_path, "r") as zf:
        names = set(zf.namelist())
        # Always keep top-level package items
        for name in names:
            if not name.startswith("ppt/"):
                kept.add(name)
        # Crawl from presentation.xml
        if "ppt/presentation.xml" in names:
            crawl(zf, "ppt/presentation.xml")
        # Always keep theme, fonts, content types
        for name in names:
            if any(x in name for x in ["theme", "fontTable", "tableStyles", "viewProps", "presProps"]):
                kept.add(name)
    return kept


def strip_orphaned_parts(pptx_path: Path):
    """
    Rewrite the PPTX zip, keeping only parts reachable from presentation.xml.
    Removes the original template slides that were left as orphaned blobs.
    Reduces file size significantly.
    """
    kept = _collect_referenced_parts(str(pptx_path))
    buf = io.BytesIO()
    with zipfile.ZipFile(str(pptx_path), "r") as zf_in, \
         zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zf_out:
        seen = set()
        for item in zf_in.infolist():
            if item.filename in kept and item.filename not in seen:
                seen.add(item.filename)
                zf_out.writestr(item, zf_in.read(item.filename))
    pptx_path.write_bytes(buf.getvalue())


# ── PPTX generator ────────────────────────────────────────────────────────────

def generate(input_json: Path, output_pptx: Path, template_pptx: Path, viz_json: Path = None):
    print(f"[pptx] Loading template: {template_pptx}")

    # Load slide JSON
    with open(input_json) as f:
        slides_data = json.load(f)
    print(f"[pptx] Slides to generate: {len(slides_data)}")

    # Load visualization specs and index by id (for native chart rendering)
    global VIZ_INDEX
    VIZ_INDEX = {}
    if viz_json and Path(viz_json).exists():
        try:
            with open(viz_json) as f:
                specs = json.load(f)
            for spec in (specs or []):
                vid = spec.get("visualization_id")
                if vid:
                    VIZ_INDEX[vid] = spec
            renderable = sum(1 for s in VIZ_INDEX.values() if normalize_chart_data(s))
            print(f"[pptx] Visualization specs: {len(VIZ_INDEX)} loaded, {renderable} renderable")
        except Exception as e:
            print(f"[pptx] Warning: could not load viz specs: {e}")

    # Open template, clear its content slides (keep masters/layouts), add ours.
    prs = Presentation(str(template_pptx))

    # Remove all existing slides from the presentation index
    xml_slides = prs.slides._sldIdLst
    rIds = [el.get(qn("r:id")) for el in list(xml_slides)]
    for el in list(xml_slides):
        xml_slides.remove(el)
    for rId in rIds:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass

    # Generate slides from JSON
    for i, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("slide_type", "unknown")
        try:
            build_slide(prs, slide_data)
            status = "ok"
        except Exception as e:
            status = f"error: {e}"
            try:
                layout = get_layout(prs, 0, 1)
                s = prs.slides.add_slide(layout)
                set_ph_text(s, 0, slide_data.get("title", f"Slide {i+1}"))
            except Exception:
                pass
        print(f"  [{i+1:2d}/{len(slides_data)}] [{slide_type:<30s}] {status}")

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    raw_kb = output_pptx.stat().st_size // 1024
    print(f"[pptx] Raw save: {raw_kb} KB — stripping orphaned template parts...")

    strip_orphaned_parts(output_pptx)
    final_kb = output_pptx.stat().st_size // 1024
    print(f"[pptx] Saved:   {output_pptx}  ({final_kb} KB)")


# ── CLI ───────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate branded PPTX from slide JSON + template")
    parser.add_argument("--input",    required=True, help="Path to slide_deck_output.json")
    parser.add_argument("--output",   required=True, help="Output .pptx path")
    parser.add_argument("--template", default=str(DEFAULT_TEMPLATE), help="Template .pptx path")
    parser.add_argument("--viz",      default=None, help="Path to visualization_specs JSON")
    args = parser.parse_args()

    generate(
        input_json   = Path(args.input),
        output_pptx  = Path(args.output),
        template_pptx= Path(args.template),
        viz_json     = Path(args.viz) if args.viz else None,
    )
